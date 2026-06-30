import uuid
import logging
import os
import io
from typing import Optional
from app.pipeline.attachments.types import AttachmentIngestResult, AttachmentChunkResult

logger = logging.getLogger(__name__)

ALLOWED_EXTENSIONS = {"pdf", "txt", "md", "pptx", "docx"}

def is_allowed_attachment(filename: str) -> bool:
    if "." not in filename:
        return False
    ext = filename.rsplit(".", 1)[1].lower()
    return ext in ALLOWED_EXTENSIONS

def _ingest_pdf(meeting_id: str, filename: str, content: bytes, attachment_id: str) -> AttachmentIngestResult:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.error("PyMuPDF (fitz) is not installed. PDF extraction failed.")
        raise ValueError("PDF extraction requires PyMuPDF to be installed.")

    doc = fitz.open(stream=content, filetype="pdf")
    page_count = len(doc)
    
    chunks = []
    
    for i in range(page_count):
        page = doc[i]
        text = page.get_text()
        if not text.strip():
            continue
            
        chunks.append(AttachmentChunkResult(
            chunk_id=str(uuid.uuid4()),
            locator_type="page",
            locator_value=str(i + 1),
            title=f"Page {i + 1}",
            text=text,
            char_count=len(text)
        ))
        
    return AttachmentIngestResult(
        attachment_id=attachment_id,
        filename=filename,
        mime_type="application/pdf",
        page_count=page_count,
        chunks=chunks,
        storage_path=f"/uploads/{meeting_id}_{filename}"
    )

def _ingest_text(meeting_id: str, filename: str, content: bytes, attachment_id: str) -> AttachmentIngestResult:
    text = content.decode("utf-8", errors="replace")
    
    # Simple chunking by paragraphs or lines for text
    # Here we just chunk by fixed length for simplicity
    chunk_size = 2000
    overlap = 200
    
    chunks = []
    i = 0
    chunk_idx = 1
    while i < len(text):
        end = min(i + chunk_size, len(text))
        chunk_text = text[i:end]
        chunks.append(AttachmentChunkResult(
            chunk_id=str(uuid.uuid4()),
            locator_type="section",
            locator_value=str(chunk_idx),
            title=f"Section {chunk_idx}",
            text=chunk_text,
            char_count=len(chunk_text)
        ))
        i += chunk_size - overlap
        chunk_idx += 1

    return AttachmentIngestResult(
        attachment_id=attachment_id,
        filename=filename,
        mime_type="text/plain",
        chunks=chunks,
        storage_path=f"/uploads/{meeting_id}_{filename}"
    )

def _ingest_pptx(meeting_id: str, filename: str, content: bytes, attachment_id: str) -> AttachmentIngestResult:
    try:
        from pptx import Presentation
        import io
    except ImportError:
        logger.error("python-pptx is not installed. PPTX extraction failed.")
        raise ValueError("PPTX extraction requires python-pptx to be installed.")

    prs = Presentation(io.BytesIO(content))
    page_count = len(prs.slides)
    
    chunks = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
                
        text = "\n".join(slide_text)
        if not text.strip():
            continue
            
        chunks.append(AttachmentChunkResult(
            chunk_id=str(uuid.uuid4()),
            locator_type="slide",
            locator_value=str(i + 1),
            title=f"Slide {i + 1}",
            text=text,
            char_count=len(text)
        ))
        
    return AttachmentIngestResult(
        attachment_id=attachment_id,
        filename=filename,
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        slide_count=page_count,
        page_count=page_count,
        chunks=chunks,
        storage_path=f"/uploads/{meeting_id}_{filename.replace('.pptx', '.pdf')}"
    )

def _ingest_docx(meeting_id: str, filename: str, content: bytes, attachment_id: str) -> AttachmentIngestResult:
    try:
        import docx
        import io
    except ImportError:
        logger.error("python-docx is not installed. DOCX extraction failed.")
        raise ValueError("DOCX extraction requires python-docx to be installed.")

    doc = docx.Document(io.BytesIO(content))
    
    chunks = []
    chunk_size = 2000
    overlap = 200
    
    full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    
    i = 0
    chunk_idx = 1
    while i < len(full_text):
        end = min(i + chunk_size, len(full_text))
        chunk_text = full_text[i:end]
        chunks.append(AttachmentChunkResult(
            chunk_id=str(uuid.uuid4()),
            locator_type="section",
            locator_value=str(chunk_idx),
            title=f"Section {chunk_idx}",
            text=chunk_text,
            char_count=len(chunk_text)
        ))
        i += chunk_size - overlap
        chunk_idx += 1

    return AttachmentIngestResult(
        attachment_id=attachment_id,
        filename=filename,
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        chunks=chunks,
        storage_path=f"/uploads/{meeting_id}_{filename.replace('.docx', '.pdf')}"
    )

def ingest_attachment(meeting_id: str, filename: str, content: bytes) -> AttachmentIngestResult:
    attachment_id = str(uuid.uuid4())
    
    ext = filename.rsplit(".", 1)[1].lower() if "." in filename else ""
    
    os.makedirs("uploads", exist_ok=True)
    out_path = os.path.join("uploads", f"{meeting_id}_{filename}")
    
    if ext == "pdf":
        with open(out_path, "wb") as f:
            f.write(content)
        return _ingest_pdf(meeting_id, filename, content, attachment_id)
    elif ext == "pptx":
        with open(out_path, "wb") as f:
            f.write(content)
        pdf_path = out_path.replace(".pptx", ".pdf")
        try:
            import comtypes.client
            comtypes.CoInitialize()
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            deck = powerpoint.Presentations.Open(os.path.abspath(out_path), WithWindow=False)
            deck.SaveAs(os.path.abspath(pdf_path), 32) # 32 = ppSaveAsPDF
            deck.Close()
            powerpoint.Quit()
            comtypes.CoUninitialize()
        except Exception as e:
            logger.error(f"Failed to convert PPTX to PDF using comtypes: {e}")
        return _ingest_pptx(meeting_id, filename, content, attachment_id)
    elif ext == "docx":
        with open(out_path, "wb") as f:
            f.write(content)
        pdf_path = out_path.replace(".docx", ".pdf")
        try:
            import comtypes.client
            comtypes.CoInitialize()
            word = comtypes.client.CreateObject("Word.Application")
            doc = word.Documents.Open(os.path.abspath(out_path))
            doc.SaveAs(os.path.abspath(pdf_path), FileFormat=17) # 17 = wdFormatPDF
            doc.Close()
            word.Quit()
            comtypes.CoUninitialize()
        except Exception as e:
            logger.error(f"Failed to convert DOCX to PDF using comtypes: {e}")
        return _ingest_docx(meeting_id, filename, content, attachment_id)
    elif ext in ["txt", "md"]:
        with open(out_path, "wb") as f:
            f.write(content)
        return _ingest_text(meeting_id, filename, content, attachment_id)
    else:
        raise ValueError(f"Unsupported file extension: {ext}")
